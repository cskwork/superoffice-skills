#!/usr/bin/env python3
"""build_demo - /superoffice 산출물 샘플을 한 번에 생성(+선택적으로 연다).

레포의 cross-platform 헬퍼 `templates/doc-env.py`와 `templates/brand-kit.example.json`
팔레트를 그대로 사용한다 - Windows/macOS/Linux에서 동일하게 동작(폰트/soffice 경로는 doc-env가 분기).

사용:
  python examples/build_demo.py            # examples/out/ 에 생성만
  python examples/build_demo.py --open     # 생성 후 기본 앱으로 연다 (OS 자동 감지)

필요: python-docx python-pptx openpyxl python-hwpx (없는 형식은 건너뛰고 사유를 출력).
PDF 변환은 LibreOffice(soffice)가 있을 때만(없으면 docx 원본만, 위조 금지).
"""
from __future__ import annotations
import importlib.util
import json
import os
import platform
import subprocess
import sys
from datetime import date

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
OUT = os.path.join(HERE, "out")
os.makedirs(OUT, exist_ok=True)

# --- 레포 헬퍼/브랜드킷 로드 (doc-env.py는 하이픈 파일명 -> spec 로더) ---
_spec = importlib.util.spec_from_file_location("doc_env", os.path.join(ROOT, "templates", "doc-env.py"))
denv = importlib.util.module_from_spec(_spec); _spec.loader.exec_module(denv)
BK = json.load(open(os.path.join(ROOT, "templates", "brand-kit.example.json"), encoding="utf-8"))

C = BK["colors"]; RAG = C["rag"]
FONT = BK["fonts"]["ko"] or denv.korean_font_name()        # 브랜드 폰트, 없으면 OS 기본
TODAY = date(2026, 6, 20).isoformat()
def nh(h: str) -> str: return h.lstrip("#")                # openpyxl/w:shd는 # 없는 hex
made: list[str] = []
skipped: list[str] = []


def make_docx():
    try:
        from docx import Document
        from docx.shared import Pt
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
    except ImportError as e:
        skipped.append(f"docx (python-docx 미설치: {e})"); return

    def shade(cell, hex_no_hash):
        shd = OxmlElement("w:shd"); shd.set(qn("w:val"), "clear"); shd.set(qn("w:fill"), hex_no_hash)
        cell._tc.get_or_add_tcPr().append(shd)

    def paint_text(cell, rgb_hex, bold=True):
        from docx.shared import RGBColor
        r = cell.paragraphs[0].runs[0] if cell.paragraphs[0].runs else cell.paragraphs[0].add_run(cell.text)
        r.font.color.rgb = RGBColor.from_string(nh(rgb_hex)); r.font.bold = bold; r.font.name = FONT
        r._element.rPr.rFonts.set(qn("w:eastAsia"), FONT)

    doc = Document()
    denv.set_docx_eastasia(doc, FONT)                      # 한글 w:eastAsia (cross-platform)
    h = doc.add_heading("2026년 1분기 실적 보고", level=1)
    for run in h.runs:                                     # 브랜드 primary 색 제목
        from docx.shared import RGBColor
        run.font.color.rgb = RGBColor.from_string(nh(C["primary"]))
    doc.add_paragraph("[결론] 1분기 매출은 계획 대비 안정적으로 운영되었다. 2분기 신규 채널 확대를 권고한다.")
    doc.add_paragraph(f"보고일 {TODAY} · {BK['company_name']} · {BK['footer']}")

    t = doc.add_table(rows=4, cols=3); t.style = "Table Grid"
    for j, htext in enumerate(["항목", "상태", "비고"]):
        t.cell(0, j).text = htext; shade(t.cell(0, j), nh(C["primary"])); paint_text(t.cell(0, j), "FFFFFF")
    body = [("일정", "green", "양호", "정상 진행"),
            ("예산", "amber", "주의", "집행 점검 필요"),
            ("리스크", "red", "경고", "완화 조치 필요")]
    for i, (k, sev, label, note) in enumerate(body, start=1):
        t.cell(i, 0).text = k; t.cell(i, 1).text = label; t.cell(i, 2).text = note
        shade(t.cell(i, 1), nh(RAG[sev]))
        paint_text(t.cell(i, 1), C["ink"] if sev == "amber" else "FFFFFF")   # amber=다크, else 흰
    p = os.path.join(OUT, "01_report.docx"); doc.save(p); made.append(p)
    return p


def make_xlsx():
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
        from openpyxl.formatting.rule import CellIsRule
        from openpyxl.worksheet.table import Table, TableStyleInfo
        from openpyxl.chart import BarChart, Reference
    except ImportError as e:
        skipped.append(f"xlsx (openpyxl 미설치: {e})"); return

    wb = Workbook(); ws = wb.active; ws.title = "KPI"
    ws["A1"] = "월"; ws["B1"] = "달성률"
    for col in ("A1", "B1"):
        ws[col].font = Font(name=FONT, bold=True, color="FFFFFF")
        ws[col].fill = PatternFill("solid", fgColor=nh(C["primary"]))
        ws[col].alignment = Alignment(horizontal="center")
    for i, (m, v) in enumerate([("1월", 0.95), ("2월", 0.72), ("3월", 0.43)], start=2):
        ws.cell(i, 1, m).font = Font(name=FONT)
        c = ws.cell(i, 2, v); c.number_format = "0.0%"
    ws["A5"] = "평균"; ws["A5"].font = Font(name=FONT, bold=True)
    ws["B5"] = "=AVERAGE(B2:B4)"; ws["B5"].number_format = "0.0%"
    # RAG 3단 조건부서식 (브랜드 색): amber=다크 텍스트
    rules = [CellIsRule(operator="greaterThanOrEqual", formula=["0.9"], fill=PatternFill("solid", fgColor=nh(RAG["green"])), font=Font(color="FFFFFF")),
             CellIsRule(operator="between", formula=["0.5", "0.9"], fill=PatternFill("solid", fgColor=nh(RAG["amber"])), font=Font(color=nh(C["ink"]))),
             CellIsRule(operator="lessThan", formula=["0.5"], fill=PatternFill("solid", fgColor=nh(RAG["red"])), font=Font(color="FFFFFF"))]
    for r in rules: ws.conditional_formatting.add("B2:B4", r)
    tab = Table(displayName="KPI표", ref="A1:B4")
    tab.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showRowStripes=True)
    ws.add_table(tab); ws.freeze_panes = "A2"
    ws.column_dimensions["A"].width = 12; ws.column_dimensions["B"].width = 14
    chart = BarChart(); chart.title = "월별 달성률"
    chart.add_data(Reference(ws, min_col=2, min_row=1, max_row=4), titles_from_data=True)
    chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=4))
    ws.add_chart(chart, "D2")
    p = os.path.join(OUT, "02_kpi.xlsx"); wb.save(p); made.append(p)
    return p


def make_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError as e:
        skipped.append(f"pptx (python-pptx 미설치: {e})"); return

    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)   # 16:9

    def band(slide, x, y, w, h, hexcol):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor.from_string(nh(hexcol))
        shp.line.fill.background(); shp.shadow.inherit = False
        return shp

    def textbox(slide, x, y, w, h, text, size, hexcol, bold=True):
        tf = slide.shapes.add_textbox(x, y, w, h).text_frame; tf.word_wrap = True
        tf.text = text; pr = tf.paragraphs[0]; pr.font.size = Pt(size); pr.font.bold = bold
        pr.font.name = FONT; pr.font.color.rgb = RGBColor.from_string(nh(hexcol))
        return tf

    # 표지: primary 풀블리드 + 흰 제목 + accent 바
    s = prs.slides.add_slide(prs.slide_layouts[6])
    band(s, 0, 0, prs.slide_width, prs.slide_height, C["primary"])
    band(s, Inches(0.8), Inches(3.4), Inches(2.2), Inches(0.12), C["accent"])
    textbox(s, Inches(0.8), Inches(2.4), Inches(11), Inches(1.2), "2026년 1분기 실적 보고", 40, "FFFFFF")
    textbox(s, Inches(0.8), Inches(3.7), Inches(11), Inches(0.6), f"{BK['company_name']} · 보고일 {TODAY}", 16, "F3F5F7", bold=False)

    # 본문: 상단 primary 띠 + 좌측 accent 바 + RAG 표
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    band(s2, 0, 0, prs.slide_width, Inches(1.0), C["primary"])
    band(s2, 0, Inches(1.0), Inches(0.15), Inches(6.5), C["accent"])
    textbox(s2, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7), "[결론] 매출 안정, 2분기 채널 확대 권고", 24, "FFFFFF")
    rows = [("항목", "상태", "비고"), ("일정", "양호", "정상"), ("예산", "주의", "점검"), ("리스크", "경고", "조치")]
    tbl = s2.shapes.add_table(4, 3, Inches(0.6), Inches(1.4), Inches(7), Inches(2.4)).table
    sev = [None, "green", "amber", "red"]
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j); cell.text = val
            run = cell.text_frame.paragraphs[0].runs[0]; run.font.name = FONT; run.font.size = Pt(14)
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor.from_string(nh(C["primary"]))
                run.font.color.rgb = RGBColor.from_string("FFFFFF"); run.font.bold = True
            elif j == 1:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor.from_string(nh(RAG[sev[i]]))
                run.font.color.rgb = RGBColor.from_string(nh(C["ink"]) if sev[i] == "amber" else "FFFFFF")
    p = os.path.join(OUT, "03_deck.pptx"); prs.save(p); made.append(p)
    return p


def make_hwpx():
    try:
        from hwpx import HwpxDocument
    except ImportError as e:
        skipped.append(f"hwpx (python-hwpx 미설치: {e})"); return
    try:
        doc = HwpxDocument.new()
        doc.add_paragraph("2026년 1분기 실적 보고")
        doc.add_paragraph(f"{BK['company_name']} · 보고일 {TODAY}")
        doc.add_paragraph("[결론] 1분기 매출은 계획 대비 안정적으로 운영되었다. 2분기 신규 채널 확대를 권고한다.")
        doc.add_table(rows=2, cols=2)
        rep = doc.validate()                              # OWPML 정합성
        p = os.path.join(OUT, "04_report.hwpx"); doc.save_to_path(p); made.append(p)
        return p
    except Exception as e:
        skipped.append(f"hwpx (생성 실패: {e})")


def make_pdf(docx_path):
    # 1순위: docx -> pdf (soffice 있을 때, 원본 서식 보존)
    if docx_path and denv.find_soffice():
        try:
            p = denv.soffice_convert(docx_path, "pdf", OUT); made.append(p); return p
        except Exception as e:
            skipped.append(f"pdf via soffice 실패: {e}")
    # 대체: HTML/CSS -> Chromium (pdf.md 1순위 경로, cross-platform, soffice 불필요)
    if denv.find_chrome():
        html = (
            "<!doctype html><html lang='ko'><head><meta charset='utf-8'><style>"
            f"@page{{size:{BK['conventions']['page_size']};margin:18mm}}"
            f"body{{font-family:'{FONT}',sans-serif;color:{C['ink']}}}"
            f".cover{{background:{C['primary']};color:#fff;padding:26px 22px;border-radius:6px}}"
            f".cover h1{{margin:0;font-size:24px}} .bar{{height:6px;width:120px;background:{C['accent']};margin:12px 0}}"
            "table{border-collapse:collapse;margin-top:18px;width:60%}"
            f"th{{background:{C['primary']};color:#fff;padding:8px;text-align:left}}"
            "td{padding:8px;border:1px solid #D0D0D0}"
            f".g{{background:{RAG['green']};color:#fff}} .a{{background:{RAG['amber']};color:{C['ink']}}} .r{{background:{RAG['red']};color:#fff}}"
            "</style></head><body>"
            f"<div class='cover'><h1>2026년 1분기 실적 보고</h1><div class='bar'></div>"
            f"<div>{BK['company_name']} · 보고일 {TODAY} · {BK['footer']}</div></div>"
            "<p>[결론] 1분기 매출은 계획 대비 안정적으로 운영되었다. 2분기 신규 채널 확대를 권고한다.</p>"
            "<table><tr><th>항목</th><th>상태</th><th>비고</th></tr>"
            "<tr><td>일정</td><td class='g'>양호</td><td>정상 진행</td></tr>"
            "<tr><td>예산</td><td class='a'>주의</td><td>집행 점검 필요</td></tr>"
            "<tr><td>리스크</td><td class='r'>경고</td><td>완화 조치 필요</td></tr></table>"
            "</body></html>"
        )
        hp = os.path.join(OUT, "_report.html"); open(hp, "w", encoding="utf-8").write(html)
        try:
            p = os.path.join(OUT, "05_report.pdf"); denv.chrome_to_pdf(hp, p); made.append(p); return p
        except Exception as e:
            skipped.append(f"pdf via chrome 실패: {e}")
    else:
        skipped.append("pdf (soffice/chrome 모두 미설치 - 위조 안 함)")


def open_file(path):
    """OS 자동 감지로 기본 앱 열기 (Windows/macOS/Linux)."""
    sysname = platform.system()
    if sysname == "Windows":
        os.startfile(path)                                # type: ignore[attr-defined]
    elif sysname == "Darwin":
        subprocess.run(["open", path], check=False)
    else:
        subprocess.run(["xdg-open", path], check=False)


def main():
    dx = make_docx(); make_xlsx(); make_pptx(); make_hwpx(); make_pdf(dx)
    print(f"=== platform: {platform.system()} · 폰트: {FONT} · soffice: {bool(denv.find_soffice())} ===")
    print(f"생성 {len(made)}건:")
    for p in made: print(f"  + {p} ({os.path.getsize(p)} bytes)")
    if skipped:
        print(f"건너뜀 {len(skipped)}건 (위조 없이 사유 기록):")
        for s in skipped: print(f"  - {s}")
    if "--open" in sys.argv:
        for p in made: open_file(p)
        print(f"열기 요청 {len(made)}건 (기본 앱).")
    return 0


if __name__ == "__main__":
    sys.exit(main())
